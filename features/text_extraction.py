import os
import PyPDF2
import pytesseract
from pdf2image import convert_from_bytes
from PIL import Image
import zipfile
from nltk.tokenize import word_tokenize
from pptx import Presentation

def zip_file_extraction(original_path, destination_path):
    with zipfile.ZipFile(original_path, 'r') as zipped:
        zipped.extractall(destination_path)

    for root, _, files in os.walk(destination_path):
        for file in files:
            file_path = os.path.join(root, file)
            print(f"Processing file: {file_path}")
            extracted_text = file_processing(file_path)

            output_file = os.path.splitext(file_path)[0] + "_extracted.txt"
            with open(output_file, "w", encoding="utf-8") as text_file:
                text_file.write(extracted_text)

def file_processing(file_path):
    file_extension = os.path.splitext(file_path)[1].lower()

    if file_extension == '.pdf':
        return extract_text_from_pdf(file_path)
    elif file_extension == '.pptx':
        return extracting_ppts(file_path)
    else:
        print(f"Unsupported file type: {file_extension}")
        return ""

def extract_text_from_pdf(file_path):
        with open(file_path, 'rb') as file:
            reader = PyPDF2.PdfReader(file)
            text = ""
            for page in reader.pages:
                text += page.extract_text()

            # If no text is extracted, perform OCR
            if not text.strip():
                print("No embedded text found. Using OCR...")
                text = perform_ocr(file_path)
            else:
                print("Embedded text extracted successfully.")

        return text

def extracting_ppts(file_path):
        content = ""
        presentation = Presentation(file_path)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    content += shape.text
        return content

def perform_ocr(file_path):
        with open(file_path, 'rb') as file:
            pdf_bytes = file.read()
        # Conversion of PDF to images
        images = convert_from_bytes(pdf_bytes)
        ocr_text = ""
        for image in images:
            ocr_text += pytesseract.image_to_string(image)

        return ocr_text

def text_preprocessing(content):
    token = word_tokenize(content.lower())
    token = [token2 for token2 in token if token2.isalnum()]
    return token


